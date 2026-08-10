#!/usr/bin/env python3
"""Generate the Lanturn Stage 1 production-readiness review deck (PPTX).

Follows the same shape as the other deck generators on this box
(Nexus/docs/generate_onboarding_deck.py): design constants at the top, small
styling helpers, one `build_slide_NN_*` per slide, `main()` assembles and saves.

Layout follows the June 2026 leadership deck it reports against — navy title
slide, white content slides under a navy header band — so the two read as one
series.

Usage:
    pip install python-pptx
    python3 scripts/generate_stage1_deck.py [output.pptx]
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ──────────────────────────────────────────────
# Design constants
# ──────────────────────────────────────────────
NAVY        = RGBColor(0x14, 0x20, 0x2C)   # title slide + header band
NAVY_RULE   = RGBColor(0x2F, 0x75, 0x95)   # accent rule under the header
ACCENT      = RGBColor(0x2F, 0x75, 0x95)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
INK         = RGBColor(0x16, 0x21, 0x2C)   # body heading
INK_2       = RGBColor(0x41, 0x52, 0x5F)   # body copy
INK_3       = RGBColor(0x6B, 0x7C, 0x88)   # captions, footer
RULE        = RGBColor(0xD9, 0xDF, 0xE3)
SURFACE_2   = RGBColor(0xEE, 0xF1, 0xF3)   # card fill
EYEBROW     = RGBColor(0x8F, 0xA7, 0xB8)   # on navy
TITLE_SUB   = RGBColor(0xB3, 0xC5, 0xD2)   # on navy

DONE        = RGBColor(0x2C, 0x6F, 0x52)
DONE_SOFT   = RGBColor(0xE2, 0xEF, 0xE8)
OPEN        = RGBColor(0x9A, 0x63, 0x18)
OPEN_SOFT   = RGBColor(0xF6, 0xEC, 0xDD)
NOTE_SOFT   = RGBColor(0xE4, 0xEE, 0xF3)

FONT      = "Calibri"
FONT_MONO = "Consolas"

SLIDE_W = Emu(12192000)   # 13.333in — 16:9
SLIDE_H = Emu(6858000)    # 7.5in

MARGIN   = Inches(0.62)
CONTENT_W = Inches(12.09)
HEADER_H = Inches(1.34)
BODY_TOP = Inches(1.72)

TOTAL_SLIDES = 9


# ──────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────
def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def create_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    set_slide_bg(slide, bg)
    return slide


def add_shape(slide, left, top, width, height, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, text, left, top, width, height,
             size=16, color=INK_2, bold=False, font=FONT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    if spacing:
        p.line_spacing = spacing
    f = p.font
    f.name, f.size, f.color.rgb, f.bold = font, Pt(size), color, bold
    return box


def add_header(slide, eyebrow, title):
    """Navy band with an uppercase eyebrow, the slide title, and an accent rule."""
    add_shape(slide, 0, 0, SLIDE_W, HEADER_H, NAVY)
    add_shape(slide, 0, HEADER_H, SLIDE_W, Pt(3), NAVY_RULE)
    add_text(slide, eyebrow.upper(), MARGIN, Inches(0.26), CONTENT_W, Inches(0.24),
             size=10.5, color=EYEBROW, font=FONT_MONO)
    add_text(slide, title, MARGIN, Inches(0.55), CONTENT_W, Inches(0.62),
             size=27, color=WHITE, bold=True)


def add_lede(slide, text, top=BODY_TOP, width=Inches(10.6)):
    box = add_text(slide, text, MARGIN, top, width, Inches(0.7),
                   size=14.5, color=INK_2, spacing=1.28)
    return box


def add_footer(slide, page):
    add_text(slide, "LANTURN — STAGE 1 REVIEW", MARGIN, Inches(6.92),
             Inches(6), Inches(0.24), size=8.5, color=INK_3, font=FONT_MONO)
    add_text(slide, f"{page:02d}", Inches(11.6), Inches(6.92), Inches(1.11),
             Inches(0.24), size=8.5, color=INK_3, font=FONT_MONO, align=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, title, body, metric=None):
    """Filled card: heading, body copy, optional accent metric line.

    The body only gives up room at the bottom when there is a metric to sit
    there — otherwise it would be squeezed into a fraction of the card and
    overflow the fill.
    """
    add_shape(slide, left, top, width, height, SURFACE_2)
    pad = Inches(0.24)
    inner_w = width - pad * 2
    body_top = top + Inches(0.56)
    reserved = Inches(0.46) if metric else Inches(0.16)
    add_text(slide, title, left + pad, top + Inches(0.2), inner_w, Inches(0.3),
             size=13.5, color=INK, bold=True)
    add_text(slide, body, left + pad, body_top, inner_w,
             height - Inches(0.56) - reserved, size=11, color=INK_2, spacing=1.3)
    if metric:
        add_text(slide, metric, left + pad, top + height - Inches(0.4),
                 inner_w, Inches(0.26), size=10.5, color=ACCENT, font=FONT_MONO)


def add_finding(slide, left, top, width, height, tag, title, body):
    """One row of the 'findings' list: mono tag on the left, text on the right."""
    tag_w = Inches(1.05)
    add_text(slide, tag, left, top + Inches(0.03), tag_w, Inches(0.26),
             size=10.5, color=ACCENT, font=FONT_MONO)
    add_text(slide, title, left + tag_w, top, width - tag_w, Inches(0.26),
             size=12.5, color=INK, bold=True)
    add_text(slide, body, left + tag_w, top + Inches(0.28), width - tag_w,
             height - Inches(0.28), size=10.5, color=INK_2, spacing=1.26)
    add_shape(slide, left, top + height + Inches(0.03), width, Pt(0.75), RULE)


def add_note(slide, top, bold_lead, body, height=Inches(0.78)):
    """Tinted callout with a 3pt accent spine."""
    add_shape(slide, MARGIN, top, CONTENT_W, height, NOTE_SOFT)
    add_shape(slide, MARGIN, top, Pt(3), height, ACCENT)
    box = slide.shapes.add_textbox(MARGIN + Inches(0.22), top + Inches(0.14),
                                   CONTENT_W - Inches(0.44), height - Inches(0.24))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.line_spacing = 1.3
    r1 = p.add_run(); r1.text = bold_lead + " "
    r1.font.name, r1.font.size, r1.font.bold, r1.font.color.rgb = FONT, Pt(11), True, INK
    r2 = p.add_run(); r2.text = body
    r2.font.name, r2.font.size, r2.font.color.rgb = FONT, Pt(11), INK_2


def style_cell(cell, text, size=10.5, color=INK_2, bold=False,
               font=FONT, fill=WHITE, anchor=MSO_ANCHOR.MIDDLE):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = anchor
    cell.margin_left = cell.margin_right = Inches(0.1)
    cell.margin_top = cell.margin_bottom = Inches(0.04)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    f = p.font
    f.name, f.size, f.color.rgb, f.bold = font, Pt(size), color, bold


# ──────────────────────────────────────────────
# Slides
# ──────────────────────────────────────────────
def build_01_title(prs):
    slide = create_slide(prs, NAVY)
    add_text(slide, "PRODUCTION READINESS · FOLLOW-UP TO 18 JUNE",
             MARGIN, Inches(2.32), CONTENT_W, Inches(0.28),
             size=11, color=EYEBROW, font=FONT_MONO)
    add_text(slide, "Stage 1 is done, bar the walkthrough",
             MARGIN, Inches(2.76), Inches(11.2), Inches(1.0),
             size=44, color=WHITE, bold=True)
    add_text(slide,
             "Eight of nine items delivered. The ninth is user acceptance "
             "testing, which needs the deployment box and the team.",
             MARGIN, Inches(3.92), Inches(9.4), Inches(0.8),
             size=16, color=TITLE_SUB, spacing=1.3)
    add_shape(slide, MARGIN, Inches(4.92), Inches(11.2), Pt(0.75),
              RGBColor(0x3A, 0x4B, 0x5A))
    add_text(slide, "Lanturn · TracePcap    ·    18 Jun – 10 Aug 2026",
             MARGIN, Inches(5.12), CONTENT_W, Inches(0.3),
             size=11, color=EYEBROW, font=FONT_MONO)
    return slide


def build_02_standing(prs):
    slide = create_slide(prs)
    add_header(slide, "The headline", "Where we stand")
    add_lede(slide,
             "In June the tool worked but was unproven at volume, unvalidated for "
             "concurrent use, and had no safety net. All three are now closed. What "
             "is left is confirming it on the real hardware with real analysts.")

    stats = [("8 / 9", "Stage 1 items delivered", DONE),
             ("110", "Substantive changes merged", INK),
             ("62", "Issues closed", INK),
             ("1", "Remaining: UAT", OPEN)]
    top, h = Inches(3.02), Inches(1.16)
    gap = Inches(0.16)
    w = (CONTENT_W - gap * 3) / 4
    for i, (value, label, colour) in enumerate(stats):
        left = MARGIN + (w + gap) * i
        add_shape(slide, left, top, w, h, SURFACE_2)
        add_text(slide, value, left + Inches(0.22), top + Inches(0.16),
                 w - Inches(0.44), Inches(0.5), size=26, color=colour,
                 bold=True, font=FONT_MONO)
        add_text(slide, label, left + Inches(0.22), top + Inches(0.74),
                 w - Inches(0.44), Inches(0.3), size=10.5, color=INK_3)

    add_note(slide, Inches(4.68), "The honest caveat.",
             "Several items are delivered as a working mechanism whose behaviour has been "
             "verified here, not yet on the production box under real load. That gap is "
             "exactly what the remaining item — UAT — exists to close, and it is written "
             "into the ticket as explicit scenarios.", height=Inches(0.92))
    add_footer(slide, 2)
    return slide


def build_03_scorecard(prs):
    slide = create_slide(prs)
    add_header(slide, "Against the June plan", "Stage 1 scorecard")

    rows = [
        ("Concurrent multi-user use",
         "Load-tested to ~20 simultaneous analyses; reads are not the constraint", "Done"),
        ("Keep large imports fast",
         "Index audit found no critical gaps; connection-pool sizing bug fixed", "Done"),
        ("Retention & capacity planning",
         "Per-file partitioning, separable packet retention, capacity reporting", "Done"),
        ("Crash / memory safety",
         "CPU and memory limits on every service, sized for native subprocesses", "Done"),
        ("Automatic backups + tested restore",
         "Nightly backup, restore rehearsed against total data loss", "Done"),
        ("Turn on production settings",
         "Production profile active; third-party AI and online geo egress closed", "Done"),
        ("Basic password hygiene",
         "Shipped defaults removed; production refuses to start without real ones", "Done"),
        ("Load test + tuning",
         "Bottleneck quantified: threat detection is ~94% of per-file cost", "Done"),
        ("UAT before go-live",
         "Scenarios written and prerequisites met; needs the box and the team", "Remaining"),
    ]

    tbl_top = Inches(1.94)
    shape = slide.shapes.add_table(len(rows) + 1, 3, MARGIN, tbl_top,
                                   CONTENT_W, Inches(4.6))
    table = shape.table
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(7.24)
    table.columns[2].width = Inches(1.35)

    headers = ["ITEM", "WHAT SHIPPED", "STATUS"]
    for c, text in enumerate(headers):
        style_cell(table.cell(0, c), text, size=9, color=INK_3,
                   bold=True, font=FONT_MONO, fill=SURFACE_2)
    table.rows[0].height = Inches(0.32)

    for r, (item, detail, status) in enumerate(rows, start=1):
        done = status == "Done"
        style_cell(table.cell(r, 0), item, size=10.5, color=INK, bold=True)
        style_cell(table.cell(r, 1), detail, size=10, color=INK_2)
        style_cell(table.cell(r, 2), status.upper(), size=9, font=FONT_MONO, bold=True,
                   color=DONE if done else OPEN,
                   fill=DONE_SOFT if done else OPEN_SOFT)
        table.rows[r].height = Inches(0.44)

    add_footer(slide, 3)
    return slide


def build_04_changed(prs):
    slide = create_slide(prs)
    add_header(slide, "In practice", "What Stage 1 actually changed")
    add_lede(slide, "Four capabilities the deployment did not have in June.")

    cards = [
        ("A safety net",
         "Nightly backup of the database, the captures and the custom rules. The restore "
         "was rehearsed by deliberately destroying everything and bringing it back — "
         "recovered data was byte-identical.", "RPO 24 h · RTO minutes"),
        ("Bounded data growth",
         "Reclaiming a capture's packets is now an instant operation rather than deleting "
         "millions of rows. Bulky packet data can expire early while the analysis "
         "summaries are kept.", "~5 ms to reclaim a file"),
        ("Predictable capacity",
         "A reporting command shows current usage and projects when the disk runs out, "
         "plus sizing guidance for choosing a retention window against a given disk.",
         "~2.5x ingest volume"),
        ("Safe startup defaults",
         "Production refuses to start on shipped passwords, with an unset AI endpoint, or "
         "without resource limits. Mistakes fail loudly at startup. Security beyond that "
         "is a separate picture — see slide 07.", "Fail fast, not fail quiet"),
    ]
    gap = Inches(0.18)
    w = (CONTENT_W - gap * 3) / 4
    for i, (title, body, metric) in enumerate(cards):
        add_card(slide, MARGIN + (w + gap) * i, Inches(2.72), w, Inches(2.98),
                 title, body, metric)
    add_footer(slide, 4)
    return slide


def build_05_uncovered(prs):
    slide = create_slide(prs)
    add_header(slide, "Worth knowing", "What the hardening work uncovered")
    add_lede(slide,
             "Each of these was live and unnoticed. None would have been found without "
             "doing the readiness work — which is the strongest argument that it was "
             "worth doing.")

    finds = [
        ("#612", "Authenticated deployment could not start from clean",
         "A permissions fault meant the login service crash-looped on any fresh install. "
         "Existing installs were unaffected, so it stayed hidden — but no new deployment "
         "could have logged in."),
        ("#628", "Retention settings did nothing at all",
         "The controls were documented and appeared to work, but never reached the "
         "application. An operator disabling deletion for evidence preservation would "
         "still have lost every capture after 12 hours."),
        ("#607", "Production defaulted to a third-party AI service",
         "On a fresh box the AI features would have sent capture-derived content to an "
         "external provider, contradicting the offline requirement. Now an explicit, "
         "deliberate setting."),
        ("#451", "Overloading the queue silently dropped work",
         "Files could be left permanently “processing” with no error surfaced. The "
         "system now applies back-pressure instead of discarding, and recovers stranded "
         "files."),
    ]
    top = Inches(2.92)
    for tag, title, body in finds:
        add_finding(slide, MARGIN, top, CONTENT_W, Inches(0.78), tag, title, body)
        top += Inches(0.92)
    add_footer(slide, 5)
    return slide


def build_06_alongside(prs):
    slide = create_slide(prs)
    add_header(slide, "18 June – 10 August", "Also delivered alongside")
    add_lede(slide,
             "Readiness was roughly a quarter of the period. The rest went into analysis "
             "capability, structure and quality.")

    themes = [
        ("Threat detection",
         "Suricata intrusion detection integrated as an offline module, with a "
         "deployment-wide switch and alerts surfaced in monitor mode."),
        ("Host identity",
         "One adjudicated verdict per host, with supporting evidence shown and "
         "disagreement rendered as contested rather than silently resolved."),
        ("Authentication",
         "Optional single sign-on, including a fully air-gapped variant with bundled "
         "identity images and a user-management guide."),
        ("Monitor mode",
         "Subnet labelling with AI assessment and staleness detection, capture timeline "
         "slicing, richer change events and security signals."),
        ("Network visualisation",
         "Node-to-node volume heatmap, protocol edge colouring with a legend, identity "
         "filtering and corrected device iconography."),
        ("Structure",
         "A five-stage architecture with enforced boundaries. Recorded violations fell "
         "from 757 to 66 and cross-module cycles to zero."),
        ("Scale & correctness",
         "An audit found lists silently truncating data; pagination moved server-side "
         "across the app so large captures render completely."),
        ("Quality gates",
         "Browser end-to-end tests, frontend unit tests, containerised integration tests "
         "and an API contract snapshot — 12 automated checks per change."),
    ]
    gap = Inches(0.18)
    w = (CONTENT_W - gap * 3) / 4
    h = Inches(1.86)
    for i, (title, body) in enumerate(themes):
        col, row = i % 4, i // 4
        add_card(slide, MARGIN + (w + gap) * col,
                 Inches(2.72) + (h + gap) * row, w, h, title, body)
    add_footer(slide, 6)
    return slide


def build_07_security(prs):
    slide = create_slide(prs)
    add_header(slide, "Security review", "Where security stands")
    add_lede(slide,
             "Stage 1 did not include a security workstream — the June plan assumed a trusted, "
             "physically-controlled network. The audit findings are therefore mostly still open, "
             "and one is worth a decision now.")

    items = [
        ("Open", "Captures are downloadable without logging in",
         "Object storage is set to public read in every deployment path, including production, and "
         "its port is published. Sign-on gates the application; it does not gate the captures "
         "behind it — so the most sensitive data is the one thing not covered. Newly raised."),
        ("Open", "No transport encryption",
         "Traffic is HTTP. Fine inside a trusted network with the browser and server on the same "
         "segment; not acceptable if the deployment is ever reached across an untrusted one."),
        ("Open", "Unauthenticated detection-rule editing",
         "Custom signature rules can be read and rewritten without credentials when the base stack "
         "runs without sign-on. Enabling the auth overlay closes this."),
        ("Partly closed", "Credentials, secrets and data egress",
         "Shipped passwords are gone and production will not start without real ones. AI and "
         "geolocation egress is closed off. Secrets are still environment variables, not a managed "
         "store — acceptable at this scale, revisit for Stage 2."),
    ]
    top = Inches(2.76)
    for tag, title, body in items:
        add_finding(slide, MARGIN, top, CONTENT_W, Inches(0.68), tag, title, body)
        top += Inches(0.80)

    add_note(slide, Inches(5.98), "Recommendation.",
             "Close the public object storage before the tool holds real operational captures — it "
             "is a small change and it is the one finding that undoes the sign-on work. The rest can "
             "be scheduled deliberately once the trust assumptions of the deployment are confirmed.",
             height=Inches(0.86))
    add_footer(slide, 7)
    return slide


def build_08_limits(prs):
    slide = create_slide(prs)
    add_header(slide, "Scope boundaries", "What Stage 1 does not cover")
    add_lede(slide,
             "Stage 1 was scoped to run today's tool safely for the team. These sit "
             "outside that line — none block the walkthrough, but each is better known "
             "now than discovered later.")

    items = [
        ("Scale", "Sized for a working window, not the whole archive",
         "Retention keeps the live data set bounded, which is what makes the single server "
         "viable. Pointing the tool at a large historical corpus is a different design — the "
         "load-test findings recommend ingesting a moving window and evicting behind it."),
        ("Ops", "No monitoring or alerting",
         "The health check works, but the metrics endpoint returns nothing — a missing "
         "dependency — so no dashboard or alert can watch the box. Today the capacity "
         "command is the only visibility, and someone has to run it."),
        ("Recovery", "A restore loses up to a day of analyst annotation",
         "Captures can be re-uploaded and re-analysed, so the packets are never truly lost. "
         "Labels, overrides and notes entered since the last nightly backup cannot be "
         "recovered — a deliberate trade, not an oversight."),
        ("Single box", "No redundancy or failover",
         "One server means hardware failure is an outage until it is repaired or rebuilt "
         "from backup. Backups protect the data; they do not keep the service running. "
         "Removing that exposure is Stage 2 work."),
    ]
    top = Inches(2.76)
    for tag, title, body in items:
        add_finding(slide, MARGIN, top, CONTENT_W, Inches(0.68), tag, title, body)
        top += Inches(0.80)

    add_note(slide, Inches(5.98), "Also open:",
             "one known defect where scheduled clean-up of monitor snapshots can leave the "
             "change history inconsistent. It is disabled by default and documented, so no "
             "deployment is exposed today.", height=Inches(0.86))
    add_footer(slide, 7)
    return slide


def build_09_recommendation(prs):
    slide = create_slide(prs)
    add_header(slide, "Next steps", "Recommendation")
    add_lede(slide,
             "Stage 1 was scoped as “safely run what we have” on a single server for "
             "roughly 10–25 analysts. That work is complete.")

    cards = [
        ("1  ·  Schedule the UAT",
         "The last Stage 1 item. Run the team's real workflow on the deployment box, with "
         "production settings and a rehearsed restore already in place. Fix what it "
         "surfaces, then sign off.", "3–5 days"),
        ("2  ·  Confirm the trade-offs",
         "Two deliberate choices are worth an explicit decision: a 24-hour backup window, "
         "and analyst annotations being the one thing a restore cannot recover.",
         "Decision, not work"),
        ("3  ·  Decide on Stage 2 separately",
         "Per-user logins and access control, shared workspaces, high availability, scaling "
         "across servers, and click-through to raw packet bytes. A product choice, not a "
         "prerequisite.", "Multi-month"),
    ]
    gap = Inches(0.2)
    w = (CONTENT_W - gap * 2) / 3
    for i, (title, body, metric) in enumerate(cards):
        add_card(slide, MARGIN + (w + gap) * i, Inches(2.86), w, Inches(2.42),
                 title, body, metric)

    add_note(slide, Inches(5.62), "Position today.",
             "The tool can be run for the team on one server with backups, bounded growth "
             "and safe defaults. It has not yet been exercised by that team on that "
             "hardware — so the recommendation is to go to UAT now, not to go live "
             "directly.", height=Inches(0.86))
    add_footer(slide, 9)
    return slide


# ──────────────────────────────────────────────
def main():
    out = sys.argv[1] if len(sys.argv) > 1 else "stage1-review-deck.pptx"

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_01_title(prs)
    build_02_standing(prs)
    build_03_scorecard(prs)
    build_04_changed(prs)
    build_05_uncovered(prs)
    build_06_alongside(prs)
    build_07_security(prs)
    build_08_limits(prs)
    build_09_recommendation(prs)

    prs.save(out)
    print(f"Saved {out}  ({len(prs.slides)} slides, expected {TOTAL_SLIDES})")


if __name__ == "__main__":
    main()
